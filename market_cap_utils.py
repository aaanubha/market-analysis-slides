# market_cap_utils.py

import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches
import datetime

# def extract_market_cap(filepath):

#     xl = pd.ExcelFile(filepath)
#     for sheet in xl.sheet_names:
#         if "Simple Model" in sheet or "Case" in sheet:
#             df = xl.parse(sheet, header=None)
#             share_count = None
#             share_price = None
            

#             # Find Share Count
#             for i, row in df.iterrows():
#                 row_str = row.astype(str).str.lower().str.strip()
#                 if row_str.str.contains("shares").any() and i + 1 < len(df):
#                     numbers = pd.to_numeric(df.iloc[i + 1].dropna(), errors='coerce').dropna()
#                     if not numbers.empty:
#                         share_count = numbers.values[0]
#                         break

#             # Find Share Price
#             for i, row in df.iterrows():
#                 row_str = row.astype(str).str.lower().str.strip()
#                 if row_str.str.contains("price").any() and i + 1 < len(df):
#                     numbers = pd.to_numeric(df.iloc[i + 1].dropna(), errors='coerce').dropna()
#                     if not numbers.empty:
#                         share_price = numbers.values[0]
#                         break

#             if share_count and share_price:
#                 market_cap = round(share_price * share_count, 2)
#                 return {
#                     "Company": os.path.basename(filepath).split(" ")[0],
#                     "Market Cap (USD)": market_cap
#                 }

#     return {"Company": os.path.basename(filepath).split(" ")[0], "Market Cap (USD)": "Error"}

import pandas as pd
import os
import re

def extract_market_cap(filepath):
    # Keywords to match share count and share price
    share_keywords = [
        "diluted shares", "fd shares", "shares outstanding", "total shares", "fully diluted",
        "shares o/s", "basic shares"
    ]
    price_keywords = [
        "share price", "price per share", "implied price", "stock price", "current price",
        "price target", "trading price"
    ]

    xl = pd.ExcelFile(filepath)
    company_name = os.path.basename(filepath).split(" ")[0]
    print(f"📄 Processing: {company_name} from {filepath}")

    for sheet_name in xl.sheet_names:
        if any(k in sheet_name.lower() for k in ["simple", "case", "model", "valuation", "base"]):
            print(f"🔍 Searching in sheet: {sheet_name}")
            df = xl.parse(sheet_name, header=None).fillna("")

            share_count = None
            share_price = None

            for i, row in df.iterrows():
                row_lower = row.astype(str).str.lower().str.strip()
                row_text = " ".join(row_lower)

                # Check for Share Count
                if any(keyword in row_text for keyword in share_keywords):
                    nums = pd.to_numeric(row, errors="coerce").dropna()
                    if not nums.empty:
                        share_count = nums.iloc[0]
                        print(f"✅ Found Share Count: {share_count} on row {i}")
                        break
                    elif i + 1 < len(df):
                        next_row_nums = pd.to_numeric(df.iloc[i + 1], errors="coerce").dropna()
                        if not next_row_nums.empty:
                            share_count = next_row_nums.iloc[0]
                            print(f"✅ Found Share Count (next row): {share_count} at row {i+1}")
                            break

            for i, row in df.iterrows():
                row_lower = row.astype(str).str.lower().str.strip()
                row_text = " ".join(row_lower)

                # Check for Share Price
                if any(keyword in row_text for keyword in price_keywords):
                    nums = pd.to_numeric(row, errors="coerce").dropna()
                    if not nums.empty:
                        share_price = nums.iloc[0]
                        print(f"✅ Found Share Price: {share_price} on row {i}")
                        break
                    elif i + 1 < len(df):
                        next_row_nums = pd.to_numeric(df.iloc[i + 1], errors="coerce").dropna()
                        if not next_row_nums.empty:
                            share_price = next_row_nums.iloc[0]
                            print(f"✅ Found Share Price (next row): {share_price} at row {i+1}")
                            break

            if share_count and share_price:
                market_cap = round(share_count * share_price, 2)
                print(f"💰 Market Cap: {market_cap}")
                return {
                    "Company": company_name,
                    "Market Cap (USD)": market_cap
                }
            else:
                print(f"❌ {company_name}: Missing value - Share Count: {share_count}, Share Price: {share_price}")

    return {
        "Company": company_name,
        "Market Cap (USD)": "Error"
    }


def create_market_cap_slide(data, pptx_path, output_path):
    
    prs = Presentation(pptx_path)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Market Capitalization Summary"

    rows, cols = len(data)+1, len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8)).table

    for j, col in enumerate(data[0].keys()):
        table.cell(0, j).text = col

    for i, row in enumerate(data):
        for j, key in enumerate(row):
            table.cell(i+1, j).text = str(row[key])

    prs.save(output_path)


# market_cap_utils.py
# market_cap_utils.py


# import pandas as pd
# import os
# from pptx import Presentation
# from pptx.util import Inches

# def extract_market_cap(filepath):
#     xl = pd.ExcelFile(filepath)
#     company = os.path.basename(filepath).split(" ")[0].upper()

#     for sheet in xl.sheet_names:
#         if "Simple Model" in sheet or "Case" in sheet:
#             df = xl.parse(sheet, header=None)
#             share_count = None
#             share_price = None

#             for i, row in df.iterrows():
#                 row_str = row.astype(str).str.lower().str.strip()
#                 if row_str.str.contains("shares").any() and i + 1 < len(df):
#                     numbers = pd.to_numeric(df.iloc[i + 1].dropna(), errors='coerce').dropna()
#                     if not numbers.empty:
#                         share_count = numbers.values[0]
#                         break

#             for i, row in df.iterrows():
#                 row_str = row.astype(str).str.lower().str.strip()
#                 if row_str.str.contains("price").any() and i + 1 < len(df):
#                     numbers = pd.to_numeric(df.iloc[i + 1].dropna(), errors='coerce').dropna()
#                     if not numbers.empty:
#                         share_price = numbers.values[0]
#                         break

#             if share_price and share_count:
#                 market_cap = round(share_price * share_count, 2)
#                 return company, market_cap

#     return company, "Error"

# def update_mkt_cap_column(pptx_path, model_folder, output_path):
#     prs = Presentation(pptx_path)

#     market_caps = {}
#     for file in os.listdir(model_folder):
#         if file.endswith(".xlsx"):
#             company, cap = extract_market_cap(os.path.join(model_folder, file))
#             market_caps[company] = cap

#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_table:
#                 table = shape.table
#                 headers = [cell.text.strip().lower() for cell in table.rows[0].cells]
#                 if "mkt. cap." in headers:
#                     mkt_col_idx = headers.index("mkt. cap.")
#                     name_col_idx = headers.index("company")

#                     for i in range(1, len(table.rows)):
#                         company = table.cell(i, name_col_idx).text.strip().upper()
#                         cap = market_caps.get(company, "Error")
#                         table.cell(i, mkt_col_idx).text = f"{cap:,}" if isinstance(cap, (int, float)) else str(cap)

#     prs.save(output_path)

